from flask import (
    Flask,
    request,
    render_template,
    redirect,
    url_for,
    jsonify,
    session,
)
from datetime import datetime
import openai
import fitz
import os
import tempfile
from dotenv import load_dotenv
from pptx import Presentation

import mysql.connector


# MySQL database configuration
db_config = {"host": "localhost", "user": "root", "password": "", "database": "eduflex"}


def connect_db():
    return mysql.connector.connect(**db_config)


load_dotenv()

app = Flask(__name__)
app.secret_key = os.getenv("APP_SECRET_KEY")
openai.api_key = os.getenv("OPENAI_API_KEY")


# Home Page
@app.route("/")
def home():
    return render_template("index.html")



# Authentication
@app.route("/login")
def login():
    return render_template("login.html")

@app.route("/login_submit", methods=["POST"])
def login_submit():
    email = request.form["email"]
    password = request.form["password"]

    # Connect to the database
    db = connect_db()
    cursor = db.cursor()

    # Query the database for user with given email and password
    query = "SELECT * FROM users WHERE email = %s AND password = %s"
    cursor.execute(query, (email, password))
    user = cursor.fetchone()
    print(user)

    if user:
        session["user"] = {
            "id": user[0],
            "nom": user[1],
            "prenom": user[2],
            "email": user[3],
            "password": user[4],
            "role": user[5],
        }
        # Determine the role based on user data
        role = user[5]

        if role == "PROFESSOR":
            return redirect(url_for("professor_index"))
        elif role == "STUDENT":
            return redirect(url_for("student_index"))

    # Invalid credentials case
    error = "Invalid credentials. Please try again."
    return render_template("login.html", error=error)



# Professor Endpoints
@app.route("/professor/index", methods=["GET", "POST"])
def professor_index():
    if request.method == "POST":
        return upload_pdf()

    # Get the current user ID from session
    user = session.get("user")
    if user is None:
        return redirect(url_for("login"))  # Ensure user is logged in
    if user["role"] != "PROFESSOR":
        return redirect(url_for("login"))

    return render_template("professor/index.html", user=user)


@app.route("/professor/feedback", methods=["GET"])
def professor_general_feedback():
    # Ensure user is logged in and is a professor
    user = session.get("user")
    if not user or user["role"] != "PROFESSOR":
        return jsonify({"error": "Unauthorized"}), 401

    # Get current user ID from session
    user_id = user["id"]

    # Connect to the database
    db = connect_db()
    cursor = db.cursor()

    try:
        # Query to fetch notes content for courses taught by the professor
        query = """
        SELECT n.content
        FROM notes n
        JOIN courses c ON n.course_id = c.id
        WHERE c.professor_id = %s
        """
        cursor.execute(query, (user_id,))
        notes = cursor.fetchall()
        notes_content = [note[0] for note in notes]  # Extract content from tuples

        # Generate general feedback in French using OpenAI API based on notes content
        if notes_content:
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[
                    {
                        "role": "system",
                        "content": "Vous êtes un assistant AI fournissant des commentaires généraux basés sur le contenu des notes académiques.",
                    },
                    {
                        "role": "user",
                        "content": f"Fournissez des commentaires généraux sur la performance des étudiants basés sur le contenu des notes suivantes:\n\n{', '.join(notes_content)}",
                    },
                ],
                temperature=0.7,
                max_tokens=1500,
            )
            general_feedback = response.choices[0].message["content"].strip()
        else:
            general_feedback = "Aucune note trouvée pour vos cours."

        cursor.close()
        db.close()

        return jsonify({"general_feedback": general_feedback})

    except Exception as e:
        return (
            jsonify(
                {
                    "error": f"Échec de la récupération du contenu des notes ou de la génération des commentaires: {str(e)}"
                }
            ),
            500,
        )


@app.route("/professor/feedbacks")
def professor_feedback():
    # Get the current user ID from session
    user = session.get("user")
    if user is None:
        return redirect(url_for("login"))  # Ensure user is logged in

    user_id = user["id"]

    # Connect to the database
    db = connect_db()
    cursor = db.cursor()

    # Fetch conversation details
    query = """
    SELECT 
        u.id, u.nom, u.prenom, u.email, 
        MAX(m.created_at) as last_message_date, 
        COUNT(m.id) as message_count
    FROM 
        messages m
    JOIN 
        users u ON (m.id_sender = u.id OR m.id_receiver = u.id)
    WHERE 
        (m.id_sender = %s OR m.id_receiver = %s) AND u.role = 'STUDENT'
    GROUP BY 
        u.id, u.nom, u.prenom, u.email
    HAVING 
        u.id != %s
    ORDER BY 
        last_message_date DESC
    """
    cursor.execute(query, (user_id, user_id, user_id))
    conversations = cursor.fetchall()

    cursor.close()
    db.close()

    # Add the default course to each conversation
    conversations_with_course = [
        {
            "id": convo[0],
            "nom": convo[1],
            "prenom": convo[2],
            "email": convo[3],
            "last_message_date": convo[4].strftime("%Y-%m-%d"),
            "message_count": convo[5],
            "course": "Gestion D'entreprise",
        }
        for convo in conversations
    ]
    number = len(conversations)

    return render_template(
        "professor/feedback.html",
        user=user,
        conversations=conversations_with_course,
        number=number,
    )


# Function to fetch subjects with professor's full name
def get_subjects():
    db = connect_db()
    cursor = db.cursor()

    # Joining courses and users to get professor's full name
    query = """
    SELECT c.title, c.description, c.id AS course_id, CONCAT(u.nom, ' ', u.prenom) AS professor_name
    FROM courses c
    JOIN users u ON c.professor_id = u.id
    """
    cursor.execute(query)
    subjects = cursor.fetchall()

    cursor.close()
    db.close()
    return subjects

@app.route("/professor/comment", methods=["GET", "POST"])
def professor_comment():
    if request.method == "POST":
        return upload_pdf()
    return render_template("professor/comment.html")


@app.route("/professor/discussion/<receiver_id>")
def professor_discussion(receiver_id):
    # Get the current user ID from session
    user = session.get("user")
    if user is None:
        return redirect(url_for("login"))  # Ensure user is logged in

    user_id = user["id"]

    # Connect to the database
    db = connect_db()
    cursor = db.cursor()

    # Fetch messages between the current user and the receiver
    query = """
    SELECT id_sender, id_receiver, message, created_at
    FROM messages
    WHERE (id_sender = %s AND id_receiver = %s)
       OR (id_sender = %s AND id_receiver = %s)
    ORDER BY created_at
    """
    cursor.execute(query, (user_id, receiver_id, receiver_id, user_id))
    messages = cursor.fetchall()

    cursor.close()
    db.close()

    return render_template(
        "professor/discussion.html",
        messages=messages,
        receiver_id=receiver_id,
        user_id=user_id,
    )



# Student Endpoints
@app.route("/student/index")
def student_index():
    user = session.get("user")
    if user is None:
        return redirect(url_for("login"))
    if user["role"] != "STUDENT":
        return redirect(url_for("login"))

    subjects = get_subjects()  # Fetch subjects from the database

    return render_template("student/index.html", user=user, subjects=subjects)




@app.route("/student/course/<course_id>", methods=["GET", "POST"])
def student_course(course_id):
    if request.method == "POST":
        # Add a new note to the database for the current user and course
        note_content = request.form.get("notes")

        # Get current user ID from session
        user = session.get("user")
        if user is None:
            return redirect(url_for("login"))  # Redirect to login if user not logged in

        user_id = user["id"]

        # Connect to the database
        db = connect_db()
        cursor = db.cursor()

        # Insert the new note into the notes table
        insert_query = (
            "INSERT INTO notes (user_id, course_id, content) VALUES (%s, %s, %s)"
        )
        cursor.execute(insert_query, (user_id, course_id, note_content))
        db.commit()

        cursor.close()
        db.close()

        # Redirect to refresh the page and show updated notes
        return redirect(url_for("student_course", course_id=course_id))

    # Fetch list of notes for the current user and course
    # This part will be executed on both GET and after POST (redirect)
    # Get current user ID from session
    user = session.get("user")
    if user is None:
        return redirect(url_for("login"))  # Redirect to login if user not logged in

    user_id = user["id"]

    # Connect to the database
    db = connect_db()
    cursor = db.cursor()

    # Fetch notes for the current user and course
    select_query = "SELECT * FROM notes WHERE user_id = %s AND course_id = %s"
    cursor.execute(select_query, (user_id, course_id))
    notes = cursor.fetchall()

    # Convert notes to list of dictionaries
    notes = [{"content": note[3], "timestamp": note[4]} for note in notes]

    # Fetch course details
    select_query_course = "SELECT * FROM courses WHERE id = %s"
    cursor.execute(select_query_course, (course_id,))
    course = cursor.fetchone()

    cursor.close()
    db.close()

    return render_template("student/course.html", notes=notes, course=course)

@app.route("/student/notes", methods=["GET"])
def student_notes():
    user = session.get("user")
    if user is None:
        return redirect(url_for("login"))  # Redirect to login if user not logged in

    user_id = user["id"]

    # Connect to the database
    db = connect_db()
    cursor = db.cursor()

    # Fetch notes for the current user
    select_query = "SELECT * FROM notes WHERE user_id = %s"
    cursor.execute(select_query, (user_id,))
    notes = cursor.fetchall()

    # Convert notes to list of dictionaries
    notes = [{"content": note[3], "timestamp": note[4]} for note in notes]
    number = len(notes)

    cursor.close()
    db.close()
    print(notes)

    return render_template("student/notes.html", notes=notes, number=number)



@app.route("/app/send/<receiver_id>", methods=["POST"])
def app_send(receiver_id):
    message = request.form.get("message")

    # Get the current user ID from session
    user = session.get("user")
    if user is None:
        return redirect(url_for("login"))  # Ensure user is logged in

    user_id = user["id"]

    # Connect to the database
    db = connect_db()
    cursor = db.cursor()

    # Insert the new message into the database
    query = "INSERT INTO messages (id_sender, id_receiver, message, created_at) VALUES (%s, %s, %s, %s)"
    cursor.execute(query, (user_id, receiver_id, message, datetime.now()))
    db.commit()

    cursor.close()
    db.close()
    
    role = user["role"]
    if role == "STUDENT":
        return redirect(url_for("student_discussion", receiver_id=receiver_id))
    else :
        return redirect(url_for("professor_discussion", receiver_id=receiver_id))

    



@app.route("/student/discussion/<receiver_id>")
def student_discussion(receiver_id):
    # Get the current user ID from session
    user = session.get("user")
    if user is None:
        return redirect(url_for("login"))  # Ensure user is logged in

    user_id = user["id"]

    # Connect to the database
    db = connect_db()
    cursor = db.cursor()

    # Fetch messages between the current user and the receiver
    query = """
    SELECT id_sender, id_receiver, message, created_at
    FROM messages
    WHERE (id_sender = %s AND id_receiver = %s)
       OR (id_sender = %s AND id_receiver = %s)
    ORDER BY created_at
    """
    cursor.execute(query, (user_id, receiver_id, receiver_id, user_id))
    messages = cursor.fetchall()

    cursor.close()
    db.close()

    return render_template(
        "student/discussion.html",
        messages=messages,
        receiver_id=receiver_id,
        user_id=user_id,
    )


# Ici pour les scripts
@app.route("/upload_pdf", methods=["POST"])
def upload_pdf():
    if "pdf_file" not in request.files:
        return jsonify({"error": "No file part"})

    file = request.files["pdf_file"]
    if file.filename == "":
        return jsonify({"error": "No selected file"})

    if file and file.filename.endswith(".pdf"):
        try:
            with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                file.save(temp_file.name)
                text = ""
                with fitz.open(temp_file.name) as doc:
                    text = "".join(page.get_text() for page in doc)

            # Split text into chunks to avoid exceeding token limits
            text_chunks = split_text_into_chunks(text, max_tokens=1500)

            generated_scripts = []
            for chunk in text_chunks:
                response = openai.ChatCompletion.create(
                    model="gpt-3.5-turbo",
                    messages=[
                        {
                            "role": "system",
                            "content": "Vous êtes un assistant qui génère des scripts détaillés pour des vidéos éducatives basées sur le contenu académique fourni. Créez un script qui est à la fois informatif et captivant, adapté pour une vidéo éducative.",
                        },
                        {
                            "role": "user",
                            "content": f"Générer un script pour une vidéo éducative basée sur le contenu suivant:\n\n{chunk}",
                        },
                    ],
                    temperature=0.7,
                    max_tokens=1500,
                )
                generated_script = response.choices[0].message["content"].strip()
                generated_scripts.append(generated_script)

            final_script = "\n\n".join(generated_scripts)
            return jsonify({"script": final_script})

        except Exception as e:
            return jsonify({"error": f"Failed to generate script: {str(e)}"})
    else:
        return jsonify({"error": "Invalid file type"})

def split_text_into_chunks(text, max_tokens):
    words = text.split()
    chunks = []
    current_chunk = []

    current_length = 0
    for word in words:
        current_length += len(word) + 1  # Add 1 for the space
        if current_length > max_tokens:
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]
            current_length = len(word) + 1
        else:
            current_chunk.append(word)

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks


@app.route("/modify_script", methods=["POST"])
def modify_script():
    try:
        # Retrieve data from the request
        script_commands = request.form.get("scriptCommands", "")
        original_script = request.form.get("originalScript", "")

        # Process script_commands as needed
        prompt = f"{original_script}\n\nCommands:\n{script_commands}"

        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {
                    "role": "system",
                    "content": "Vous êtes un assistant qui génère des scripts détaillés pour des vidéos éducatives basées sur le contenu académique fourni. Créez un script qui est à la fois informatif et captivant, adapté pour une vidéo éducative.",
                },
                {
                    "role": "user",
                    "content": prompt,
                },
            ],
            temperature=0.7,
            max_tokens=1500,
        )
        generated_script = response.choices[0].message["content"].strip()

        # Return updated script as JSON response
        return jsonify(updatedScript=generated_script)

    except Exception as e:
        return jsonify({"error": f"Failed to modify script: {str(e)}"})


# Ici pour les commentaires
@app.route("/upload_file", methods=["POST"])
def upload_file():
    if "pdf_file" not in request.files:
        return jsonify({"error": "No file part"})

    file = request.files["pdf_file"]
    if file.filename == "":
        return jsonify({"error": "No selected file"})

    if file and file.filename.endswith(".pdf"):
        try:
            if file and allowed_file(file.filename):
                extension = file.filename.rsplit(".", 1)[1].lower()
                if extension == "pptx":
                    texts = extract_text_from_ppt(file)
                elif extension == "pdf":
                    texts = extract_text_from_pdf(file)
                comments = generate_comments(texts)
                return jsonify({"script": comments})

        except Exception as e:
            return jsonify({"error": f"Failed to generate comments: {str(e)}"})
    else:
        return jsonify({"error": "Invalid file type"})


@app.route("/modify_comment", methods=["POST"])
def modify_comment():
    try:
        # Retrieve data from the request
        script_commands = request.form.get("scriptCommands", "")
        original_script = request.form.get("originalScript", "")

        # Process script_commands as needed
        prompt = f"{original_script}\n\nCommands:\n{script_commands}"

        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {
                    "role": "system",
                    "content": "Vous êtes un assistant qui génère des scripts détaillés pour des vidéos éducatives basées sur le contenu académique fourni. Créez un script qui est à la fois informatif et captivant, adapté pour une vidéo éducative.",
                },
                {
                    "role": "user",
                    "content": prompt,
                },
            ],
            temperature=0.7,
            max_tokens=1500,
        )
        generated_script = response.choices[0].message["content"].strip()

        # Return updated script as JSON response
        return jsonify(updatedScript=generated_script)

    except Exception as e:
        return jsonify({"error": f"Failed to modify script: {str(e)}"})


def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in {"pptx", "pdf"}


def extract_text_from_ppt(file):
    prs = Presentation(file)
    return [extract_text(slide) for slide in prs.slides]


def extract_text(slide):
    return " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])


def extract_text_from_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    return [page.get_text() for page in doc]


def generate_comments(texts):
    comments = []
    for text in texts:
        try:
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",  # Confirm this model name from OpenAI documentation
                messages=[
                    {
                        "role": "system",
                        "content": "Vous êtes un assistant qui fournit des commentaires détaillés sur le contenu académique.",
                    },
                    {
                        "role": "user",
                        "content": f"Fournissez des commentaires perspicaces pour chaque page ou diapositive de ce cours. Les commentaires doivent être pertinents par rapport au contenu de la page:\n\n{text}",
                    },
                ],
                temperature=0.7,
                max_tokens=1500,
            )
            comments.append(response.choices[0].message["content"].strip())
        except Exception as e:
            comments.append(f"Failed to generate comment: {str(e)}")
    return comments


if __name__ == "__main__":
    app.run(debug=True)
